# -*- coding: utf-8 -*-
"""
Created on Thu Feb 27 15:36:51 2025

@author: SaiAnuhyaKurra
"""
import os
import re
from pptx import Presentation
from pptx.util import Inches
from collections import defaultdict

def get_project_images(input_dir):
    """Fetch all image files from the input directory and categorize them project-wise."""
    project_images = defaultdict(lambda: {"map": [], "legend": [], "percent": []})
    
    # Updated regex pattern to allow flexible project names
    pattern = re.compile(r"(.+?)(_leg_map|_map_percent|_map|_legend|_percent)?\.png", re.IGNORECASE)
    
    for file in sorted(os.listdir(input_dir)):
        if file.lower().endswith(('.png', '.jpg', '.jpeg')):
            match = pattern.match(file)
            if match:
                project_name, image_type = match.groups()
                file_path = os.path.join(input_dir, file)

                # Handle cases where no suffix is found
                if image_type is None:
                    project_images[project_name]["map"].append(file_path)
                elif image_type in ["_leg_map", "_legend"]:
                    project_images[project_name]["legend"].append(file_path)
                elif image_type in ["_map_percent", "_percent"]:
                    project_images[project_name]["percent"].append(file_path)
                else:
                    project_images[project_name]["map"].append(file_path)

    return project_images

def create_presentation(input_dir, output_ppt):
    """Generate a PowerPoint presentation with images grouped by project."""
    project_images = get_project_images(input_dir)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define order of projects
    project_order = ["WS", "NCF", "WR", "Elevation", "Solar", "Terrain"]
    
    sorted_projects = sorted(project_images.keys(), key=lambda x: next((i for i, p in enumerate(project_order) if p in x), len(project_order)))
    
    for project in sorted_projects:
        images = project_images[project]
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Remove default title placeholder
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide.shapes._spTree.remove(shape._element)
        
        # Add Main Map (Left Side)
        if images["map"]:
            slide.shapes.add_picture(images["map"][0], Inches(0.1), Inches(0.87), width=Inches(6.45), height=Inches(6.3))
        
        # Add Legend (Inside Bottom-Left of Map)
        if images["legend"]:
            slide.shapes.add_picture(images["legend"][0], Inches(0.1), Inches(5.71), width=Inches(1.09), height=Inches(1.46))
        
        # Add Percentage Map (Bottom-Right)
        if images["percent"]:
            slide.shapes.add_picture(images["percent"][0], Inches(6.67), Inches(4.63), width=Inches(3), height=Inches(2))
        
        # Add a description textbox (Above Percentage Map)
        text_box = slide.shapes.add_textbox(Inches(6.67), Inches(2.42), Inches(3.16), Inches(0.71))
        text_frame = text_box.text_frame
        text_frame.text = f"{project}"
    
    prs.save(output_ppt)
    print(f"Presentation saved: {output_ppt}")

# Example Usage
input_folder =r"C:\Users\SaiAnuhyaKurra\Downloads\Terrain complexity_legend"
output_pptx = "D:/Test/PPT_Task/dynamic_presentation_2.pptx"
create_presentation(input_folder, output_pptx)
#%%
#Version_1
# import os
# import re
# from pptx import Presentation
# from pptx.util import Inches
# from collections import defaultdict

# def get_project_images(input_dir):
#     """Fetch all image files from the input directory and categorize them project-wise."""
#     project_images = defaultdict(lambda: {"map": [], "legend": [], "percent": []})
    
#     pattern = re.compile(r"(.*?)(_leg_map|_map_percent|_map|_legend|_percent)\.png", re.IGNORECASE)
    
#     for file in sorted(os.listdir(input_dir)):
#         if file.lower().endswith(('.png', '.jpg', '.jpeg')):
#             match = pattern.match(file)
#             if match:
#                 project_name, image_type = match.groups()
#                 file_path = os.path.join(input_dir, file)
                
#                 if image_type == "_leg_map" or image_type == "_legend" :
#                     project_images[project_name]["legend"].append(file_path)
#                 elif image_type == "_map_percent" or  image_type == "_percent": 
#                     project_images[project_name]["percent"].append(file_path)
#                 else:
#                     project_images[project_name]["map"].append(file_path)
    
#     return project_images

# def create_presentation(input_dir, output_ppt):
#     """Generate a PowerPoint presentation with images grouped by project."""
#     project_images = get_project_images(input_dir)
#     prs = Presentation()
#     prs.slide_width = Inches(10)
#     prs.slide_height = Inches(7.5)
    
#     # Define order of projects
#     project_order = ["WS", "NCF", "WR", "Elevation", "Solar", "Terrain"]
    
#     sorted_projects = sorted(project_images.keys(), key=lambda x: next((i for i, p in enumerate(project_order) if p in x), len(project_order)))
    
#     for project in sorted_projects:
#         images = project_images[project]
#         slide = prs.slides.add_slide(prs.slide_layouts[5])
        
#         # Remove default title placeholder
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 slide.shapes._spTree.remove(shape._element)
        
#         # Add Main Map (Left Side)
#         if images["map"]:
#             slide.shapes.add_picture(images["map"][0], Inches(0.1), Inches(0.87), width=Inches(6.45), height=Inches(6.3))
        
#         # Add Legend (Inside Bottom-Left of Map)
#         if images["legend"]:
#             slide.shapes.add_picture(images["legend"][0], Inches(0.1), Inches(5.71), width=Inches(1.09), height=Inches(1.46))
        
#         # Add Percentage Map (Bottom-Right)
#         if images["percent"]:
#             slide.shapes.add_picture(images["percent"][0], Inches(6.67), Inches(4.63), width=Inches(3), height=Inches(2))
        
#         # Add a description textbox (Above Percentage Map)
#         text_box = slide.shapes.add_textbox(Inches(6.67), Inches(2.42), Inches(3.16), Inches(0.71))
#         text_frame = text_box.text_frame
#         text_frame.text = f"{project}"
    
#     prs.save(output_ppt)
#     print(f"Presentation saved: {output_ppt}")

# # Example Usage
# input_folder =r"C:\Users\SaiAnuhyaKurra\Downloads\Terrain complexity_legend"
# output_pptx = "D:/Test/PPT_Task/dynamic_presentation.pptx"
# create_presentation(input_folder, output_pptx)